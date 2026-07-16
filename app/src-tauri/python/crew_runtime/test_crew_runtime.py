from __future__ import annotations

import json
import os
import subprocess
import sys
import tempfile
import threading
import unittest
from pathlib import Path
from unittest import mock


RUNTIME_DIR = Path(__file__).resolve().parent
if str(RUNTIME_DIR) not in sys.path:
    sys.path.insert(0, str(RUNTIME_DIR))

import crew_tools
import main as crew_runtime

TEST_OPENROUTER_NEMOTRON_MODEL = "nvidia/nemotron-3-nano-30b-a3b:free"
TEST_OPENROUTER_COMPLEX_NEMOTRON_MODEL = "nvidia/nemotron-3-ultra-550b-a55b:free"


class CrewRuntimeStatusTests(unittest.TestCase):
    def test_status_requires_the_pinned_runtime_and_office_dependencies(self) -> None:
        status = crew_runtime.runtime_status()

        self.assertTrue(status["runtimeCompatible"])
        self.assertTrue(status["toolDependenciesInstalled"])
        self.assertEqual(status["crewaiVersion"], crew_runtime.EXPECTED_CREWAI_VERSION)
        self.assertEqual(status["runtimeSchemaVersion"], crew_runtime.RUNTIME_SCHEMA_VERSION)


class CrewRuntimeToolTests(unittest.TestCase):
    def setUp(self) -> None:
        self.temporary = tempfile.TemporaryDirectory()
        self.root = Path(self.temporary.name)
        self.request = {
            "cwd": str(self.root),
            "config": {
                "baseUrl": "http://127.0.0.1:11434",
                "model": "qwen3:4b",
                "timeoutMs": 30_000,
            },
            "providerConfigs": {
                "openRouter": {
                    "baseUrl": "https://openrouter.ai/api/v1",
                    "model": TEST_OPENROUTER_NEMOTRON_MODEL,
                    "apiKey": "test-key",
                    "timeoutMs": 30_000,
                }
            },
            "governance": {
                "agentAccess": [
                    {
                        "agentId": "agent-test",
                        "allowedTools": [
                            "read_file",
                            "edit_file",
                            "create_directory",
                            "glob",
                            "grep",
                            "web_fetch",
                            "web_search",
                            "bash",
                            "office_workflow",
                        ],
                        "blockedTools": ["copy_path"],
                    }
                ]
            },
        }
        self.agent = {
            "id": "agent-test",
            "name": "Test Agent",
            "role": "executor",
            "goal": "Verify the runtime",
            "backstory": "A deterministic test agent.",
            "providerKind": "openrouter",
            "tools": [
                "read_file",
                "edit_file",
                "create_directory",
                "copy_path",
                "glob",
                "grep",
                "web_fetch",
                "web_search",
                "bash",
                "office_workflow",
            ],
            "maxIterations": 2,
        }

    def tearDown(self) -> None:
        self.temporary.cleanup()

    def _tools(self) -> dict[str, object]:
        return {
            tool.name: tool
            for tool in crew_tools.build_runtime_tools(self.request, self.agent)
        }

    def test_governance_binds_allowed_tools_and_omits_blocked_tools(self) -> None:
        tools = self._tools()

        self.assertIn("web_search", tools)
        self.assertIn("office_workflow", tools)
        self.assertNotIn("copy_path", tools)

        built_agent = crew_runtime.build_agent(self.request, self.agent)
        self.assertEqual(set(tools), {tool.name for tool in built_agent.tools})

    def test_file_tools_edit_read_glob_and_grep_inside_workspace(self) -> None:
        tools = self._tools()

        write_result = tools["edit_file"]._run("src/example.py", content="print('alpha')\n")
        read_result = tools["read_file"]._run("src/example.py")
        grep_result = tools["grep"]._run("alpha", path="src", file_pattern="*.py")
        glob_result = tools["glob"]._run("**/*.py")
        escape_result = tools["edit_file"]._run("../escape.txt", content="blocked")

        self.assertIn("Updated", write_result)
        self.assertIn("print('alpha')", read_result)
        self.assertIn("src/example.py:1", grep_result)
        self.assertIn("src/example.py", glob_result)
        self.assertIn("outside the authorized working directory", escape_result)
        self.assertFalse((self.root.parent / "escape.txt").exists())

    def test_bash_uses_runtime_python_without_inheriting_pythonhome(self) -> None:
        with mock.patch.dict(
            os.environ,
            {
                "PYTHONHOME": str(self.root / "incompatible-python-home"),
                "PYTHONPATH": str(self.root / "incompatible-python-path"),
                "OPENROUTER_API_KEY": "must-not-reach-child",
            },
            clear=False,
        ):
            environment = crew_tools._subprocess_environment()
            result = self._tools()["bash"]._run(
                "python -c \"import json; print(json.dumps({'verified': True}))\"",
                30,
            )

        self.assertNotIn("PYTHONHOME", environment)
        self.assertNotIn("PYTHONPATH", environment)
        self.assertNotIn("OPENROUTER_API_KEY", environment)
        self.assertEqual(Path(environment["PATH"].split(os.pathsep)[0]), Path(sys.executable).resolve().parent)
        self.assertIn("Exit code: 0", result)
        self.assertIn('"verified": true', result)

    def test_web_fetch_blocks_private_network_destinations(self) -> None:
        result = self._tools()["web_fetch"]._run("http://127.0.0.1:11434/api/tags")

        self.assertIn("Private, loopback", result)

    def test_web_fetch_safely_extracts_oversized_html_instead_of_failing(self) -> None:
        class FakeHeaders:
            @staticmethod
            def get_content_type() -> str:
                return "text/html"

            @staticmethod
            def get_content_charset() -> str:
                return "utf-8"

        class FakeResponse:
            headers = FakeHeaders()
            status = 200

            def __enter__(self):
                return self

            def __exit__(self, *_args) -> None:
                return None

            @staticmethod
            def read(_amount: int) -> bytes:
                return ("<html><body><h1>CrewAI docs</h1>" + "useful research " * 80_000).encode("utf-8")

            @staticmethod
            def geturl() -> str:
                return "https://example.com/docs"

        fake_opener = mock.Mock()
        fake_opener.open.return_value = FakeResponse()
        with (
            mock.patch.object(crew_tools, "_validate_public_url", side_effect=lambda value: value),
            mock.patch.object(crew_tools.urllib.request, "build_opener", return_value=fake_opener),
        ):
            result = self._tools()["web_fetch"]._run("https://example.com/docs")

        self.assertIn("Download truncated safely", result)
        self.assertIn("CrewAI docs", result)
        self.assertNotIn("ERROR", result)

    def test_office_workflow_creates_a_real_powerpoint(self) -> None:
        sections = json.dumps([
            {"title": "Research", "bullets": ["Search works", "Sources included"]},
            {"title": "Coding", "body": "Files can be edited and verified."},
        ])

        result = self._tools()["office_workflow"]._run(
            "artifacts/runtime-proof.pptx",
            "Crew runtime proof",
            sections,
        )

        from pptx import Presentation

        output = self.root / "artifacts" / "runtime-proof.pptx"
        presentation = Presentation(output)
        self.assertIn("Created", result)
        self.assertTrue(output.is_file())
        self.assertGreater(output.stat().st_size, 10_000)
        self.assertEqual(len(presentation.slides), 3)


class CrewRuntimeTaskTests(unittest.TestCase):
    def test_openrouter_respects_selected_models_and_agent_overrides(self) -> None:
        request = {
            "providerConfigs": {
                "openRouter": {
                    "model": "anthropic/claude-sonnet-4",
                    "apiKey": "test-key",
                }
            }
        }

        configured = crew_runtime.resolve_agent_model_label(
            request,
            {"providerKind": "openrouter"},
        )
        overridden = crew_runtime.resolve_agent_model_label(
            request,
            {"providerKind": "openrouter", "modelOverride": "google/gemini-2.5-pro"},
        )

        self.assertEqual(configured, "anthropic/claude-sonnet-4")
        self.assertEqual(overridden, "google/gemini-2.5-pro")

    def test_openrouter_configuration_requires_an_api_key(self) -> None:
        payload = {
            "providerConfigs": {"openRouter": {"model": TEST_OPENROUTER_NEMOTRON_MODEL}},
            "maxParallelTasks": 3,
        }
        agents = [{"providerKind": "openrouter"}]

        with self.assertRaisesRegex(ValueError, "OpenRouter API key is missing"):
            crew_runtime.validate_runtime_provider_models(payload, agents)

    def test_free_openrouter_crews_are_bounded_to_two_parallel_tasks(self) -> None:
        payload = {
            "providerConfigs": {
                "openRouter": {
                    "model": TEST_OPENROUTER_NEMOTRON_MODEL,
                    "apiKey": "test-key",
                }
            },
            "maxParallelTasks": 3,
        }
        agents = [{"providerKind": "openrouter"}]

        crew_runtime.validate_runtime_provider_models(payload, agents)

        self.assertEqual(payload["maxParallelTasks"], 2)
        self.assertIn("bounded parallel execution", payload["_rateLimitPolicyReason"])

    def test_external_provider_timeout_and_free_retry_policy_are_applied(self) -> None:
        request = {
            "retryCount": 0,
            "config": {"timeoutMs": 60_000},
            "providerConfigs": {
                "openRouter": {
                    "model": TEST_OPENROUTER_NEMOTRON_MODEL,
                    "apiKey": "test-key",
                    "timeoutMs": 180_000,
                }
            },
        }
        agent = {
            "id": "runtime-agent",
            "name": "Runtime Agent",
            "role": "executor",
            "goal": "Verify provider settings.",
            "providerKind": "openrouter",
            "tools": [],
        }

        llm = crew_runtime.build_llm(request, agent)
        built_agent = crew_runtime.build_agent(request, agent)

        self.assertEqual(crew_runtime.resolve_agent_timeout_ms(request, agent), 180_000)
        self.assertEqual(llm.timeout, 180)
        self.assertTrue(llm.is_litellm)
        self.assertTrue(llm.supports_function_calling())
        self.assertEqual(llm.additional_params.get("max_retries"), 4)
        self.assertIs(llm.additional_params.get("_skip_mcp_handler"), True)
        self.assertEqual(built_agent.max_execution_time, 180)

    def test_configured_retry_count_can_raise_provider_default(self) -> None:
        request = {"retryCount": 5}

        self.assertEqual(
            crew_runtime.resolve_llm_max_retries(request, "openrouter/paid-model"),
            5,
        )

    def test_recovered_provider_limit_is_reported_as_warning_after_success(self) -> None:
        logs = [{
            "action": "runtime_stderr",
            "result": "ResourceExhausted: Worker local total request limit reached (32/32)",
        }]

        crew_runtime.mark_recovered_provider_logs(logs, "completed")

        self.assertEqual(logs[0]["action"], "provider_retry_recovered")
        self.assertEqual(logs[0]["severity"], "warning")
        self.assertEqual(logs[0]["phase"], "status")

    def test_textual_named_tool_call_is_bridged_to_native_crewai_shape(self) -> None:
        tools = [{
            "type": "function",
            "function": {
                "name": "edit_file",
                "parameters": {
                    "type": "object",
                    "properties": {"path": {"type": "string"}, "content": {"type": "string"}},
                    "required": ["path"],
                },
            },
        }]

        bridged = crew_runtime.bridge_textual_tool_call(
            '<tool_call>{"name":"edit_file","arguments":{"path":"artifacts/proof.py","content":"ok"}}</tool_call>',
            tools,
        )

        self.assertIsInstance(bridged, list)
        self.assertEqual(bridged[0]["function"]["name"], "edit_file")
        self.assertEqual(
            json.loads(bridged[0]["function"]["arguments"]),
            {"path": "artifacts/proof.py", "content": "ok"},
        )

    def test_textual_bare_arguments_are_bridged_only_for_a_unique_schema(self) -> None:
        tools = [
            {
                "type": "function",
                "function": {
                    "name": "read_file",
                    "parameters": {
                        "type": "object",
                        "properties": {"path": {"type": "string"}},
                        "required": ["path"],
                    },
                },
            },
            {
                "type": "function",
                "function": {
                    "name": "edit_file",
                    "parameters": {
                        "type": "object",
                        "properties": {"path": {"type": "string"}, "content": {"type": "string"}},
                        "required": ["path"],
                    },
                },
            },
        ]

        bridged = crew_runtime.bridge_textual_tool_call(
            '{"path":"artifacts/proof.py","content":"print(1)"}',
            tools,
        )
        ambiguous = crew_runtime.bridge_textual_tool_call('{"path":"artifacts/proof.py"}', tools)

        self.assertEqual(bridged[0]["function"]["name"], "edit_file")
        self.assertEqual(ambiguous, '{"path":"artifacts/proof.py"}')

    def test_normal_json_answer_is_not_mistaken_for_a_tool_call(self) -> None:
        tools = [{
            "type": "function",
            "function": {
                "name": "bash",
                "parameters": {
                    "type": "object",
                    "properties": {"command": {"type": "string"}},
                    "required": ["command"],
                },
            },
        }]
        answer = '{"workflow":"complex-crew","verified":true}'

        self.assertEqual(crew_runtime.bridge_textual_tool_call(answer, tools), answer)

    def test_textual_tool_sequence_bridges_only_the_first_dependent_call(self) -> None:
        tools = [
            {
                "type": "function",
                "function": {
                    "name": "edit_file",
                    "parameters": {
                        "type": "object",
                        "properties": {"path": {"type": "string"}, "content": {"type": "string"}},
                        "required": ["path"],
                    },
                },
            },
            {
                "type": "function",
                "function": {
                    "name": "bash",
                    "parameters": {
                        "type": "object",
                        "properties": {"command": {"type": "string"}},
                        "required": ["command"],
                    },
                },
            },
        ]

        bridged = crew_runtime.bridge_textual_tool_call(
            '{"path":"artifacts/proof.py","content":"print(1)"}\n\n'
            '{"command":"python artifacts/proof.py"}',
            tools,
        )

        self.assertEqual(len(bridged), 1)
        self.assertEqual(bridged[0]["function"]["name"], "edit_file")

    def test_required_office_artifact_must_exist_inside_workspace(self) -> None:
        with tempfile.TemporaryDirectory() as temporary:
            root = Path(temporary)
            request = {"cwd": str(root)}
            agent = {"tools": ["office_workflow"]}
            task = {
                "description": "Create artifacts/verified-output.pptx with office_workflow.",
                "expectedOutput": "The exact path to artifacts/verified-output.pptx.",
            }

            expected = crew_runtime.expected_office_artifact_paths(request, task, agent)

            self.assertEqual(expected, [root / "artifacts" / "verified-output.pptx"])
            self.assertEqual(crew_runtime.missing_required_artifacts(request, task, agent), expected)

            expected[0].parent.mkdir(parents=True)
            expected[0].write_bytes(b"valid-artifact")
            self.assertEqual(crew_runtime.missing_required_artifacts(request, task, agent), [])

    def test_required_edit_artifact_must_exist_inside_workspace(self) -> None:
        with tempfile.TemporaryDirectory() as temporary:
            root = Path(temporary)
            request = {"cwd": str(root)}
            agent = {"tools": ["edit_file", "bash"]}
            task = {
                "description": "Create artifacts/verified-script.py and execute it.",
                "expectedOutput": "The verified-script.py path and command output.",
            }

            expected = crew_runtime.expected_edit_artifact_paths(request, task, agent)

            self.assertEqual(expected, [root / "artifacts" / "verified-script.py"])
            self.assertEqual(crew_runtime.missing_required_artifacts(request, task, agent), expected)

            expected[0].parent.mkdir(parents=True)
            expected[0].write_text("print('verified')\n", encoding="utf-8")
            self.assertEqual(crew_runtime.missing_required_artifacts(request, task, agent), [])

    def test_deterministic_office_fallback_builds_slides_from_markdown(self) -> None:
        with tempfile.TemporaryDirectory() as temporary:
            root = Path(temporary)
            artifacts = root / "artifacts"
            artifacts.mkdir()
            (artifacts / "report.md").write_text(
                "# Evidence\n- Source A\n# Implementation Proof\n- JSON verified\n"
                "# Collaboration\n- Handoffs verified\n# Limitations\n- Free concurrency is capped at two\n",
                encoding="utf-8",
            )
            request = {"cwd": str(root), "name": "Fallback Test"}
            agent = {"tools": ["office_workflow"]}
            task = {
                "description": (
                    "Read artifacts/report.md and create artifacts/fallback.pptx "
                    "with title 'Verified Fallback'."
                ),
                "expectedOutput": "The exact path to artifacts/fallback.pptx.",
            }
            missing = crew_runtime.missing_required_artifacts(request, task, agent)

            result = crew_runtime.deterministic_office_fallback(request, task, agent, missing)

            from pptx import Presentation

            output = artifacts / "fallback.pptx"
            self.assertIsNotNone(result)
            self.assertTrue(output.is_file())
            self.assertEqual(len(Presentation(output).slides), 5)

    def test_tasks_are_stably_topologically_ordered(self) -> None:
        tasks = [
            {"id": "review", "context": ["implement"], "dependencies": []},
            {"id": "plan", "context": [], "dependencies": []},
            {"id": "implement", "context": [], "dependencies": ["plan"]},
        ]

        ordered = crew_runtime.order_task_payloads(tasks)

        self.assertEqual([task["id"] for task in ordered], ["plan", "implement", "review"])

    def test_parallel_batches_group_independent_tasks(self) -> None:
        tasks = [{"id": f"task-{index}"} for index in range(5)]

        process, normalized = crew_runtime.normalize_task_concurrency(tasks, "parallel", 2)

        self.assertEqual(process, "parallel")
        self.assertEqual(
            [task["asyncExecution"] for task in normalized],
            [True, True, True, True, False],
        )
        self.assertEqual(
            [task["_parallelBatch"] for task in normalized],
            [0, 0, 1, 1, 2],
        )

    def test_parallel_batches_wait_for_dependencies(self) -> None:
        tasks = [
            {"id": "research", "dependencies": []},
            {"id": "code", "dependencies": []},
            {"id": "synthesis", "dependencies": ["research", "code"]},
            {"id": "verify", "dependencies": ["synthesis"]},
            {"id": "present", "dependencies": ["synthesis"]},
        ]

        _, normalized = crew_runtime.normalize_task_concurrency(tasks, "parallel", 2)

        self.assertEqual(
            [(task["id"], task["_parallelBatch"]) for task in normalized],
            [("research", 0), ("code", 0), ("synthesis", 1), ("verify", 2), ("present", 2)],
        )

    def test_single_task_parallel_run_is_synchronous(self) -> None:
        _, normalized = crew_runtime.normalize_task_concurrency([{"id": "only"}], "parallel", 4)

        self.assertFalse(normalized[0]["asyncExecution"])


@unittest.skipUnless(
    os.environ.get("OPEN_COWORK_RUN_CREW_INTEGRATION") == "1"
    and bool(os.environ.get("OPENROUTER_API_KEY", "").strip()),
    "Set OPEN_COWORK_RUN_CREW_INTEGRATION=1 and OPENROUTER_API_KEY to run the live OpenRouter smoke test.",
)
class CrewRuntimeIntegrationTests(unittest.TestCase):
    def test_live_research_coding_and_powerpoint_run(self) -> None:
        with tempfile.TemporaryDirectory() as temporary:
            root = Path(temporary)
            selected_task = os.environ.get("OPEN_COWORK_CREW_TEST_TASK", "all").strip().lower()
            agents = [
                {
                    "id": "researcher",
                    "name": "Researcher",
                    "role": "researcher",
                    "goal": "Find verifiable sources.",
                    "backstory": "A careful web researcher.",
                    "providerKind": "openrouter",
                    "tools": ["web_search", "web_fetch"],
                    "allowDelegation": False,
                    "maxIterations": 3,
                },
                {
                    "id": "coder",
                    "name": "Coder",
                    "role": "executor",
                    "goal": "Create and verify working code.",
                    "backstory": "A precise software engineer.",
                    "providerKind": "openrouter",
                    "tools": ["read_file", "edit_file", "bash"],
                    "allowDelegation": False,
                    "maxIterations": 4,
                },
                {
                    "id": "presenter",
                    "name": "Presenter",
                    "role": "writer",
                    "goal": "Create real presentation artifacts.",
                    "backstory": "A concise presentation author.",
                    "providerKind": "openrouter",
                    "tools": ["office_workflow"],
                    "allowDelegation": False,
                    "maxIterations": 4,
                },
            ]
            tasks = [
                {
                    "id": "research",
                    "description": "Use web_search for 'CrewAI official documentation'. Return at least two source URLs from the tool result.",
                    "expectedOutput": "At least two real source URLs.",
                    "agentId": "researcher",
                    "context": [],
                    "dependencies": [],
                    "asyncExecution": False,
                },
                {
                    "id": "code",
                    "description": "Use edit_file to create artifacts/proof.py containing exactly print('crew tools work'), then use bash to execute it.",
                    "expectedOutput": "The created file path and successful command output.",
                    "agentId": "coder",
                    "context": [],
                    "dependencies": [],
                    "asyncExecution": False,
                },
                {
                    "id": "presentation",
                    "description": "Use office_workflow to create artifacts/proof.pptx with title 'Crew Tools Work' and two content slides.",
                    "expectedOutput": "The exact path of a valid PPTX artifact.",
                    "agentId": "presenter",
                    "context": [],
                    "dependencies": [],
                    "asyncExecution": False,
                },
            ]
            if selected_task != "all":
                tasks = [task for task in tasks if task["id"] == selected_task]
                if not tasks:
                    self.fail(f"Unknown OPEN_COWORK_CREW_TEST_TASK: {selected_task}")
            used_agent_ids = {task["agentId"] for task in tasks}
            agents = [agent for agent in agents if agent["id"] in used_agent_ids]
            allowed_by_agent = {agent["id"]: agent["tools"] for agent in agents}
            payload = {
                "id": "integration-smoke",
                "name": "Integration Smoke Crew",
                "description": "Verify research, coding, and PowerPoint tools.",
                "executionGuidelines": "Call the required tool; do not answer from memory. Finish immediately after the required tool result is verified.",
                "knowledgeFocus": "Runtime verification",
                "outputMode": "standard",
                "retryCount": 0,
                "process": "sequential",
                "maxParallelTasks": 1,
                "maxRpm": 0,
                "verbose": False,
                "cwd": str(root),
                "config": {
                    "baseUrl": "http://127.0.0.1:11434",
                    "model": "unused-for-openrouter",
                    "timeoutMs": 60_000,
                },
                "providerConfigs": {
                    "openRouter": {
                        "baseUrl": "https://openrouter.ai/api/v1",
                        "model": os.environ.get(
                            "OPEN_COWORK_CREW_TEST_MODEL",
                            TEST_OPENROUTER_NEMOTRON_MODEL,
                        ),
                        "apiKey": os.environ["OPENROUTER_API_KEY"],
                        "timeoutMs": 180_000,
                        "verifyTlsCertificates": True,
                    }
                },
                "agents": agents,
                "tasks": tasks,
                "governance": {
                    "subject": "integration-test",
                    "subjectRoles": ["owner"],
                    "pendingApprovalTypes": [],
                    "agentAccess": [
                        {
                            "agentId": agent_id,
                            "allowedTools": tools,
                            "blockedTools": [],
                            "allowedMcpServerNames": [],
                            "blockedMcpServerNames": [],
                            "delegationAllowed": False,
                            "gatewayHints": [],
                        }
                        for agent_id, tools in allowed_by_agent.items()
                    ],
                },
                "memoryContext": {},
            }

            response = crew_runtime.execute_definition(payload)

            self.assertEqual(response["status"], "completed", response.get("error"))
            if selected_task in {"all", "code"}:
                self.assertTrue((root / "artifacts" / "proof.py").is_file())
                self.assertEqual(
                    (root / "artifacts" / "proof.py").read_text(encoding="utf-8").strip(),
                    "print('crew tools work')",
                )
            if selected_task in {"all", "presentation"}:
                presentation = root / "artifacts" / "proof.pptx"
                self.assertTrue(presentation.is_file())
                self.assertGreater(presentation.stat().st_size, 10_000)
            if selected_task in {"all", "research"}:
                research_output = next(item["output"] for item in response["taskResults"] if item["taskId"] == "research")
                self.assertGreaterEqual(str(research_output).count("https://"), 2)


@unittest.skipUnless(
    os.environ.get("OPEN_COWORK_RUN_PARALLEL_CREW_INTEGRATION") == "1"
    and bool(os.environ.get("OPENROUTER_API_KEY", "").strip()),
    "Set OPEN_COWORK_RUN_PARALLEL_CREW_INTEGRATION=1 and OPENROUTER_API_KEY to run the live parallel Crew test.",
)
class CrewRuntimeParallelIntegrationTests(unittest.TestCase):
    def test_two_free_model_agents_start_llm_requests_concurrently(self) -> None:
        with tempfile.TemporaryDirectory() as temporary:
            root = Path(temporary)
            barrier = threading.Barrier(2)
            lock = threading.Lock()
            active_calls = 0
            max_active_calls = 0
            first_calls: set[int] = set()
            llm_class = crew_runtime.runtime_llm_class()
            original_llm_call = llm_class.call

            def synchronized_llm_call(llm, *args, **kwargs):
                nonlocal active_calls, max_active_calls
                with lock:
                    llm_id = id(llm)
                    is_first_call = llm_id not in first_calls
                    first_calls.add(llm_id)
                    if is_first_call:
                        active_calls += 1
                        max_active_calls = max(max_active_calls, active_calls)
                if not is_first_call:
                    return original_llm_call(llm, *args, **kwargs)
                try:
                    barrier.wait(timeout=30)
                    return original_llm_call(llm, *args, **kwargs)
                finally:
                    with lock:
                        active_calls -= 1

            agents = [
                {
                    "id": "probe-a",
                    "name": "Parallel Probe A",
                    "role": "executor",
                    "goal": "Execute the exact probe command.",
                    "backstory": "A minimal concurrency probe agent.",
                    "providerKind": "openrouter",
                    "tools": ["bash"],
                    "allowDelegation": False,
                    "maxIterations": 3,
                },
                {
                    "id": "probe-b",
                    "name": "Parallel Probe B",
                    "role": "executor",
                    "goal": "Execute the exact probe command.",
                    "backstory": "A minimal concurrency probe agent.",
                    "providerKind": "openrouter",
                    "tools": ["bash"],
                    "allowDelegation": False,
                    "maxIterations": 3,
                },
            ]
            tasks = [
                {
                    "id": "probe-a-task",
                    "description": "Call bash exactly once with command python -c \"print('PARALLEL_A')\" and return its output.",
                    "expectedOutput": "PARALLEL_A from the bash result.",
                    "agentId": "probe-a",
                    "context": [],
                    "dependencies": [],
                    "asyncExecution": False,
                },
                {
                    "id": "probe-b-task",
                    "description": "Call bash exactly once with command python -c \"print('PARALLEL_B')\" and return its output.",
                    "expectedOutput": "PARALLEL_B from the bash result.",
                    "agentId": "probe-b",
                    "context": [],
                    "dependencies": [],
                    "asyncExecution": False,
                },
            ]
            payload = {
                "id": "parallel-integration-smoke",
                "name": "Parallel Integration Crew",
                "executionGuidelines": "Call bash immediately and finish after its successful result.",
                "outputMode": "standard",
                "retryCount": 1,
                "process": "parallel",
                "maxParallelTasks": 2,
                "maxRpm": 0,
                "verbose": False,
                "cwd": str(root),
                "config": {
                    "baseUrl": "http://127.0.0.1:11434",
                    "model": "unused-for-openrouter",
                    "timeoutMs": 60_000,
                },
                "providerConfigs": {
                    "openRouter": {
                        "baseUrl": "https://openrouter.ai/api/v1",
                        "model": TEST_OPENROUTER_NEMOTRON_MODEL,
                        "apiKey": os.environ["OPENROUTER_API_KEY"],
                        "timeoutMs": 300_000,
                        "verifyTlsCertificates": True,
                    }
                },
                "agents": agents,
                "tasks": tasks,
                "governance": {
                    "subject": "parallel-integration-test",
                    "subjectRoles": ["owner"],
                    "pendingApprovalTypes": [],
                    "agentAccess": [
                        {
                            "agentId": agent["id"],
                            "allowedTools": ["bash"],
                            "blockedTools": [],
                            "allowedMcpServerNames": [],
                            "blockedMcpServerNames": [],
                            "delegationAllowed": False,
                            "gatewayHints": [],
                        }
                        for agent in agents
                    ],
                },
                "memoryContext": {},
            }

            with mock.patch.object(llm_class, "call", synchronized_llm_call):
                response = crew_runtime.execute_definition(payload)

            self.assertEqual(response["status"], "completed", response.get("error"))
            self.assertEqual(max_active_calls, 2, "Both LLM requests must be active at the same time.")
            self.assertFalse(barrier.broken)
            outputs = {result["taskId"]: str(result.get("output") or "") for result in response["taskResults"]}
            self.assertIn("PARALLEL_A", outputs["probe-a-task"])
            self.assertIn("PARALLEL_B", outputs["probe-b-task"])
            self.assertTrue(any(
                log.get("action") == "rate_limit_policy"
                and "max 2" in str(log.get("result") or "")
                for log in response["logs"]
            ))
            first_handoff = next(
                log for log in response["logs"]
                if log.get("action") == "task_handoff" and log.get("taskId") == "probe-a-task"
            )
            self.assertIn("Async: yes", str(first_handoff.get("result") or ""))


@unittest.skipUnless(
    os.environ.get("OPEN_COWORK_RUN_COMPLEX_CREW_INTEGRATION") == "1"
    and bool(os.environ.get("OPENROUTER_API_KEY", "").strip()),
    "Set OPEN_COWORK_RUN_COMPLEX_CREW_INTEGRATION=1 and OPENROUTER_API_KEY to run the complex live Crew test.",
)
class CrewRuntimeComplexIntegrationTests(unittest.TestCase):
    def test_live_complex_dependency_workflow(self) -> None:
        with tempfile.TemporaryDirectory() as temporary:
            root = Path(temporary)
            agents = [
                {
                    "id": "researcher",
                    "name": "Evidence Researcher",
                    "role": "researcher",
                    "goal": "Collect current, verifiable primary sources.",
                    "backstory": "A careful web researcher who never invents citations.",
                    "providerKind": "openrouter",
                    "modelOverride": TEST_OPENROUTER_NEMOTRON_MODEL,
                    "tools": ["web_search", "web_fetch"],
                    "allowDelegation": False,
                    "maxIterations": 4,
                },
                {
                    "id": "coder",
                    "name": "Implementation Engineer",
                    "role": "executor",
                    "goal": "Create and execute a machine-verifiable implementation proof.",
                    "backstory": "A precise engineer who verifies every generated file.",
                    "providerKind": "openrouter",
                    "tools": ["create_directory", "edit_file", "read_file", "bash"],
                    "allowDelegation": False,
                    "maxIterations": 4,
                },
                {
                    "id": "analyst",
                    "name": "Synthesis Analyst",
                    "role": "analyst",
                    "goal": "Combine research and implementation evidence into one report.",
                    "backstory": "A structured analyst who preserves upstream evidence.",
                    "providerKind": "openrouter",
                    "tools": ["edit_file", "read_file"],
                    "allowDelegation": False,
                    "maxIterations": 4,
                },
                {
                    "id": "verifier",
                    "name": "Quality Verifier",
                    "role": "reviewer",
                    "goal": "Verify files, commands, citations, and required report sections.",
                    "backstory": "A skeptical reviewer who reports concrete evidence.",
                    "providerKind": "openrouter",
                    "tools": ["read_file", "grep", "bash"],
                    "allowDelegation": False,
                    "maxIterations": 4,
                },
                {
                    "id": "presenter",
                    "name": "Presentation Author",
                    "role": "writer",
                    "goal": "Turn verified evidence into a real PowerPoint artifact.",
                    "backstory": "A concise presentation author who uses only verified inputs.",
                    "providerKind": "openrouter",
                    "tools": ["read_file", "office_workflow"],
                    "allowDelegation": False,
                    "maxIterations": 4,
                },
            ]
            tasks = [
                {
                    "id": "research",
                    "description": (
                        "Use web_search for 'CrewAI official documentation agents tasks processes' and web_fetch at least "
                        "one official result. Return a concise evidence brief with at least three real https URLs, source "
                        "titles, and key facts. Do not answer from memory."
                    ),
                    "expectedOutput": "An evidence brief containing at least three real source URLs.",
                    "agentId": "researcher",
                    "context": [],
                    "dependencies": [],
                    "asyncExecution": True,
                },
                {
                    "id": "code",
                    "description": (
                        "Use create_directory and edit_file to create artifacts/metrics.py. The script must print valid "
                        "JSON with workflow='complex-crew', agents=5, tasks=5, and verified=true. Execute it with bash and "
                        "return the exact file path plus command output."
                    ),
                    "expectedOutput": "The metrics.py path and its successful JSON output.",
                    "agentId": "coder",
                    "context": [],
                    "dependencies": [],
                    "asyncExecution": True,
                },
                {
                    "id": "synthesis",
                    "description": (
                        "Use both upstream task results. Create artifacts/complex-report.md with edit_file. It must contain "
                        "the headings Evidence, Implementation Proof, Collaboration, and Limitations; at least three URLs "
                        "from the research result; and the implementation JSON facts. Read the file back and return its path."
                    ),
                    "expectedOutput": "The verified path to a complete Markdown synthesis report.",
                    "agentId": "analyst",
                    "context": ["research", "code"],
                    "dependencies": ["research", "code"],
                    "asyncExecution": False,
                },
                {
                    "id": "verification",
                    "description": (
                        "Read artifacts/complex-report.md and artifacts/metrics.py. Execute metrics.py with bash and use grep "
                        "to verify the report headings and URLs. Return VERIFICATION_OK only with concrete command evidence."
                    ),
                    "expectedOutput": "VERIFICATION_OK followed by command and file evidence.",
                    "agentId": "verifier",
                    "context": ["synthesis"],
                    "dependencies": ["synthesis"],
                    "asyncExecution": False,
                },
                {
                    "id": "presentation",
                    "description": (
                        "Read artifacts/complex-report.md, then use office_workflow to create artifacts/complex-crew.pptx "
                        "with title 'Complex Crew Verification' and at least four content slides covering Evidence, "
                        "Implementation Proof, Collaboration, and Limitations. Return the exact created path."
                    ),
                    "expectedOutput": "The exact path to a valid PowerPoint file with at least five total slides.",
                    "agentId": "presenter",
                    "context": ["synthesis"],
                    "dependencies": ["synthesis"],
                    "asyncExecution": False,
                },
            ]
            allowed_by_agent = {agent["id"]: agent["tools"] for agent in agents}
            payload = {
                "id": "complex-integration-smoke",
                "name": "Complex Integration Crew",
                "description": "Research, implement, synthesize, verify, and present one evidence-bound result.",
                "executionGuidelines": (
                    "Use every required tool, preserve upstream evidence in downstream tasks, and never claim success "
                    "without reading or executing the created artifact."
                ),
                "knowledgeFocus": "CrewAI runtime collaboration verification",
                "outputMode": "standard",
                "retryCount": 1,
                "process": "parallel",
                "maxParallelTasks": 2,
                "maxRpm": 0,
                "verbose": False,
                "cwd": str(root),
                "config": {
                    "baseUrl": "http://127.0.0.1:11434",
                    "model": "unused-for-openrouter",
                    "timeoutMs": 60_000,
                },
                "providerConfigs": {
                    "openRouter": {
                        "baseUrl": "https://openrouter.ai/api/v1",
                        "model": os.environ.get(
                            "OPEN_COWORK_CREW_TEST_MODEL",
                            TEST_OPENROUTER_COMPLEX_NEMOTRON_MODEL,
                        ),
                        "apiKey": os.environ["OPENROUTER_API_KEY"],
                        "timeoutMs": 300_000,
                        "verifyTlsCertificates": True,
                    }
                },
                "agents": agents,
                "tasks": tasks,
                "governance": {
                    "subject": "complex-integration-test",
                    "subjectRoles": ["owner"],
                    "pendingApprovalTypes": [],
                    "agentAccess": [
                        {
                            "agentId": agent_id,
                            "allowedTools": tools,
                            "blockedTools": [],
                            "allowedMcpServerNames": [],
                            "blockedMcpServerNames": [],
                            "delegationAllowed": False,
                            "gatewayHints": [],
                        }
                        for agent_id, tools in allowed_by_agent.items()
                    ],
                },
                "memoryContext": {},
            }

            response = crew_runtime.execute_definition(payload)

            self.assertEqual(response["status"], "completed", response.get("error"))
            self.assertEqual(
                [result["taskId"] for result in response["taskResults"]],
                ["research", "code", "synthesis", "verification", "presentation"],
            )
            self.assertTrue(all(result["status"] == "completed" for result in response["taskResults"]))
            self.assertTrue(any(log.get("action") == "rate_limit_policy" for log in response["logs"]))
            self.assertEqual(
                {log.get("agentId") for log in response["logs"] if log.get("action") == "task_handoff"},
                {"researcher", "coder", "analyst", "verifier", "presenter"},
            )
            parallel_handoffs = {
                log.get("taskId")
                for log in response["logs"]
                if log.get("action") == "task_handoff"
                and "Async: yes" in str(log.get("result") or "")
            }
            self.assertEqual(parallel_handoffs, {"research", "code", "verification", "presentation"})
            outputs = {result["taskId"]: str(result.get("output") or "") for result in response["taskResults"]}
            self.assertGreaterEqual(outputs["research"].count("https://"), 3)
            self.assertIn("verified", outputs["code"].lower())
            self.assertNotIn("exit code: 1", outputs["code"].lower())
            self.assertNotIn("execution failed", outputs["code"].lower())
            self.assertIn("VERIFICATION_OK", outputs["verification"])

            metrics_path = root / "artifacts" / "metrics.py"
            report_path = root / "artifacts" / "complex-report.md"
            presentation_path = root / "artifacts" / "complex-crew.pptx"
            self.assertTrue(metrics_path.is_file())
            metrics_run = subprocess.run(
                [sys.executable, str(metrics_path)],
                cwd=root,
                capture_output=True,
                text=True,
                timeout=30,
                check=True,
            )
            metrics = json.loads(metrics_run.stdout)
            self.assertEqual(
                metrics,
                {"workflow": "complex-crew", "agents": 5, "tasks": 5, "verified": True},
            )

            report = report_path.read_text(encoding="utf-8")
            for heading in ("Evidence", "Implementation Proof", "Collaboration", "Limitations"):
                self.assertIn(heading, report)
            self.assertGreaterEqual(report.count("https://"), 3)

            self.assertTrue(presentation_path.is_file())
            self.assertGreater(presentation_path.stat().st_size, 10_000)
            from pptx import Presentation

            presentation = Presentation(presentation_path)
            self.assertGreaterEqual(len(presentation.slides), 5)


if __name__ == "__main__":
    unittest.main()
