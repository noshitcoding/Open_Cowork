from __future__ import annotations

import fnmatch
import base64
import html
import ipaddress
import json
import os
import re
import shutil
import socket
import subprocess
import sys
import tempfile
import urllib.parse
import urllib.request
from html.parser import HTMLParser
from pathlib import Path
from typing import Any, Literal

from crewai.tools import BaseTool  # type: ignore
from pydantic import BaseModel, Field, PrivateAttr


MAX_FILE_READ_BYTES = 250_000
MAX_FILE_WRITE_CHARS = 1_500_000
MAX_WEB_BYTES = 1_000_000
MAX_TOOL_OUTPUT_CHARS = 24_000
IGNORED_DIRECTORY_NAMES = {".git", ".next", ".venv", "dist", "node_modules", "target"}


def _workspace_root(request: dict) -> Path:
    configured = str(request.get("cwd") or "").strip()
    candidate = Path(configured).expanduser() if configured else Path.cwd()
    try:
        resolved = candidate.resolve(strict=True)
    except (OSError, RuntimeError):
        resolved = Path.cwd().resolve()
    return resolved if resolved.is_dir() else resolved.parent


def _resolve_workspace_path(root: Path, value: str, *, allow_root: bool = True) -> Path:
    raw = str(value or "").strip()
    if not raw:
        target = root
    else:
        candidate = Path(raw).expanduser()
        target = (candidate if candidate.is_absolute() else root / candidate).resolve(strict=False)
    try:
        target.relative_to(root)
    except ValueError as exc:
        raise ValueError(f"Path is outside the authorized working directory: {value}") from exc
    if not allow_root and target == root:
        raise ValueError("The working-directory root itself cannot be modified.")
    return target


def _truncate(value: object, limit: int = MAX_TOOL_OUTPUT_CHARS) -> str:
    text = str(value or "")
    if len(text) <= limit:
        return text
    return text[:limit].rstrip() + f"\n...[truncated after {limit} characters]"


def _safe_result(operation: str, callback) -> str:
    try:
        return _truncate(callback())
    except Exception as exc:
        return f"ERROR ({operation}): {exc.__class__.__name__}: {exc}"


def _agent_access(request: dict, agent_id: str) -> dict:
    governance = request.get("governance") or {}
    for entry in governance.get("agentAccess") or []:
        if isinstance(entry, dict) and str(entry.get("agentId") or "").strip() == agent_id:
            return entry
    return {}


def _canonical_tool_id(value: str) -> str:
    normalized = str(value or "").strip().lower().replace("-", "_")
    aliases = {
        "shell": "bash",
        "bashtool": "bash",
        "read": "read_file",
        "filereadtool": "read_file",
        "edit": "edit_file",
        "write": "edit_file",
        "fileedittool": "edit_file",
        "webfetch": "web_fetch",
        "websearch": "web_search",
        "mcp_call": "mcp",
        "generate_office_workflow": "office_workflow",
        "pptx_template_workflow": "office_workflow",
        "docx_template_workflow": "office_workflow",
    }
    return aliases.get(normalized, normalized)


class ReadFileInput(BaseModel):
    path: str = Field(description="Workspace-relative or absolute path inside the working directory")
    start_line: int = Field(default=1, ge=1, description="First 1-based line to return")
    max_lines: int = Field(default=400, ge=1, le=2000, description="Maximum number of lines")


class ReadFileTool(BaseTool):
    name: str = "read_file"
    description: str = "Read a UTF-8 text file inside the authorized working directory with line numbers."
    args_schema: type[BaseModel] = ReadFileInput
    _root: Path = PrivateAttr()

    def __init__(self, root: Path) -> None:
        super().__init__()
        self._root = root

    def _run(self, path: str, start_line: int = 1, max_lines: int = 400) -> str:
        def execute() -> str:
            target = _resolve_workspace_path(self._root, path)
            if not target.is_file():
                raise FileNotFoundError(f"File not found: {target}")
            if target.stat().st_size > MAX_FILE_READ_BYTES:
                raise ValueError(f"File exceeds the {MAX_FILE_READ_BYTES}-byte read limit")
            raw = target.read_bytes()
            if b"\x00" in raw:
                raise ValueError("Binary files are not supported by read_file")
            text = raw.decode("utf-8", errors="replace")
            lines = text.splitlines()
            start = max(0, start_line - 1)
            selected = lines[start:start + max_lines]
            rendered = "\n".join(f"{index + start + 1}: {line}" for index, line in enumerate(selected))
            return f"File: {target}\nLines: {start + 1}-{start + len(selected)} of {len(lines)}\n\n{rendered}"

        return _safe_result("read_file", execute)


class EditFileInput(BaseModel):
    path: str = Field(description="File path inside the working directory")
    content: str = Field(default="", description="Complete new file content; use this or old_text/new_text")
    old_text: str = Field(default="", description="Exact existing text to replace")
    new_text: str = Field(default="", description="Replacement text used with old_text")
    replace_all: bool = Field(default=False, description="Replace every occurrence of old_text")


class EditFileTool(BaseTool):
    name: str = "edit_file"
    description: str = "Create or edit a UTF-8 text file atomically. Use content for a full write, or old_text/new_text for a precise replacement."
    args_schema: type[BaseModel] = EditFileInput
    _root: Path = PrivateAttr()

    def __init__(self, root: Path) -> None:
        super().__init__()
        self._root = root

    def _run(
        self,
        path: str,
        content: str = "",
        old_text: str = "",
        new_text: str = "",
        replace_all: bool = False,
    ) -> str:
        def execute() -> str:
            target = _resolve_workspace_path(self._root, path, allow_root=False)
            if target.exists() and not target.is_file():
                raise ValueError(f"Target is not a file: {target}")
            if old_text:
                if not target.is_file():
                    raise FileNotFoundError(f"Cannot replace text in missing file: {target}")
                current = target.read_text(encoding="utf-8", errors="strict")
                occurrences = current.count(old_text)
                if occurrences == 0:
                    raise ValueError("old_text was not found; re-read the file before editing")
                if occurrences > 1 and not replace_all:
                    raise ValueError(f"old_text occurs {occurrences} times; provide a unique match or set replace_all")
                updated = current.replace(old_text, new_text, -1 if replace_all else 1)
                change = f"replaced {occurrences if replace_all else 1} occurrence(s)"
            else:
                updated = content
                change = "wrote complete content"
            if len(updated) > MAX_FILE_WRITE_CHARS:
                raise ValueError(f"Content exceeds the {MAX_FILE_WRITE_CHARS}-character write limit")
            target.parent.mkdir(parents=True, exist_ok=True)
            with tempfile.NamedTemporaryFile("w", encoding="utf-8", newline="", dir=target.parent, delete=False) as handle:
                handle.write(updated)
                temporary = Path(handle.name)
            os.replace(temporary, target)
            return f"Updated {target} ({change}, {len(updated)} characters)."

        return _safe_result("edit_file", execute)


class PathInput(BaseModel):
    path: str = Field(description="Path inside the working directory")


class CreateDirectoryTool(BaseTool):
    name: str = "create_directory"
    description: str = "Create a directory, including missing parent directories, inside the working directory."
    args_schema: type[BaseModel] = PathInput
    _root: Path = PrivateAttr()

    def __init__(self, root: Path) -> None:
        super().__init__()
        self._root = root

    def _run(self, path: str) -> str:
        return _safe_result("create_directory", lambda: self._create(path))

    def _create(self, path: str) -> str:
        target = _resolve_workspace_path(self._root, path, allow_root=False)
        target.mkdir(parents=True, exist_ok=True)
        return f"Created directory: {target}"


class TransferPathInput(BaseModel):
    source: str = Field(description="Existing source path inside the working directory")
    destination: str = Field(description="Destination path inside the working directory")


class MovePathTool(BaseTool):
    name: str = "move_path"
    description: str = "Move or rename a file or directory within the working directory."
    args_schema: type[BaseModel] = TransferPathInput
    _root: Path = PrivateAttr()

    def __init__(self, root: Path) -> None:
        super().__init__()
        self._root = root

    def _run(self, source: str, destination: str) -> str:
        def execute() -> str:
            src = _resolve_workspace_path(self._root, source, allow_root=False)
            dst = _resolve_workspace_path(self._root, destination, allow_root=False)
            if not src.exists():
                raise FileNotFoundError(f"Source does not exist: {src}")
            if dst.exists():
                raise FileExistsError(f"Destination already exists: {dst}")
            dst.parent.mkdir(parents=True, exist_ok=True)
            shutil.move(str(src), str(dst))
            return f"Moved {src} to {dst}"

        return _safe_result("move_path", execute)


class CopyPathTool(BaseTool):
    name: str = "copy_path"
    description: str = "Copy a file or directory within the working directory without overwriting an existing destination."
    args_schema: type[BaseModel] = TransferPathInput
    _root: Path = PrivateAttr()

    def __init__(self, root: Path) -> None:
        super().__init__()
        self._root = root

    def _run(self, source: str, destination: str) -> str:
        def execute() -> str:
            src = _resolve_workspace_path(self._root, source, allow_root=False)
            dst = _resolve_workspace_path(self._root, destination, allow_root=False)
            if not src.exists():
                raise FileNotFoundError(f"Source does not exist: {src}")
            if dst.exists():
                raise FileExistsError(f"Destination already exists: {dst}")
            dst.parent.mkdir(parents=True, exist_ok=True)
            if src.is_dir():
                shutil.copytree(src, dst)
            else:
                shutil.copy2(src, dst)
            return f"Copied {src} to {dst}"

        return _safe_result("copy_path", execute)


class GlobInput(BaseModel):
    pattern: str = Field(description="Glob pattern such as **/*.py")
    path: str = Field(default=".", description="Directory inside the working directory")
    max_results: int = Field(default=200, ge=1, le=1000)


class GlobTool(BaseTool):
    name: str = "glob"
    description: str = "Find workspace files by glob pattern."
    args_schema: type[BaseModel] = GlobInput
    _root: Path = PrivateAttr()

    def __init__(self, root: Path) -> None:
        super().__init__()
        self._root = root

    def _run(self, pattern: str, path: str = ".", max_results: int = 200) -> str:
        def execute() -> str:
            base = _resolve_workspace_path(self._root, path)
            if not base.is_dir():
                raise NotADirectoryError(f"Not a directory: {base}")
            matches: list[str] = []
            for candidate in base.glob(pattern):
                relative_parts = candidate.relative_to(self._root).parts
                if any(part in IGNORED_DIRECTORY_NAMES for part in relative_parts):
                    continue
                matches.append(candidate.relative_to(self._root).as_posix() + ("/" if candidate.is_dir() else ""))
                if len(matches) >= max_results:
                    break
            return f"Found {len(matches)} path(s):\n" + "\n".join(matches)

        return _safe_result("glob", execute)


class GrepInput(BaseModel):
    pattern: str = Field(description="Regular expression or literal text to search for")
    path: str = Field(default=".", description="File or directory inside the working directory")
    file_pattern: str = Field(default="*", description="Filename glob, for example *.py")
    case_sensitive: bool = False
    max_results: int = Field(default=200, ge=1, le=1000)


class GrepTool(BaseTool):
    name: str = "grep"
    description: str = "Search UTF-8 workspace files and return path, line number, and matching line."
    args_schema: type[BaseModel] = GrepInput
    _root: Path = PrivateAttr()

    def __init__(self, root: Path) -> None:
        super().__init__()
        self._root = root

    def _run(
        self,
        pattern: str,
        path: str = ".",
        file_pattern: str = "*",
        case_sensitive: bool = False,
        max_results: int = 200,
    ) -> str:
        def execute() -> str:
            target = _resolve_workspace_path(self._root, path)
            regex = re.compile(pattern, 0 if case_sensitive else re.IGNORECASE)
            candidates = [target] if target.is_file() else target.rglob("*")
            matches: list[str] = []
            for candidate in candidates:
                if not candidate.is_file() or not fnmatch.fnmatch(candidate.name, file_pattern):
                    continue
                relative_parts = candidate.relative_to(self._root).parts
                if any(part in IGNORED_DIRECTORY_NAMES for part in relative_parts):
                    continue
                try:
                    if candidate.stat().st_size > MAX_FILE_READ_BYTES:
                        continue
                    raw = candidate.read_bytes()
                    if b"\x00" in raw:
                        continue
                    lines = raw.decode("utf-8", errors="replace").splitlines()
                except OSError:
                    continue
                for line_number, line in enumerate(lines, 1):
                    if regex.search(line):
                        matches.append(f"{candidate.relative_to(self._root).as_posix()}:{line_number}: {_truncate(line, 500)}")
                        if len(matches) >= max_results:
                            return f"Found at least {len(matches)} match(es):\n" + "\n".join(matches)
            return f"Found {len(matches)} match(es):\n" + "\n".join(matches)

        return _safe_result("grep", execute)


class _TextExtractor(HTMLParser):
    def __init__(self) -> None:
        super().__init__(convert_charrefs=True)
        self.parts: list[str] = []
        self.hidden_depth = 0

    def handle_starttag(self, tag: str, attrs: list[tuple[str, str | None]]) -> None:
        if tag in {"script", "style", "noscript", "svg"}:
            self.hidden_depth += 1
        elif tag in {"br", "p", "div", "li", "h1", "h2", "h3", "h4", "tr"}:
            self.parts.append("\n")

    def handle_endtag(self, tag: str) -> None:
        if tag in {"script", "style", "noscript", "svg"} and self.hidden_depth:
            self.hidden_depth -= 1
        elif tag in {"p", "div", "li", "h1", "h2", "h3", "h4", "tr"}:
            self.parts.append("\n")

    def handle_data(self, data: str) -> None:
        if not self.hidden_depth:
            self.parts.append(data)

    def text(self) -> str:
        value = html.unescape("".join(self.parts)).replace("\r", "")
        lines = [" ".join(line.split()) for line in value.split("\n")]
        return "\n".join(line for line in lines if line)


def _validate_public_url(value: str) -> str:
    parsed = urllib.parse.urlsplit(str(value or "").strip())
    if parsed.scheme not in {"http", "https"} or not parsed.hostname:
        raise ValueError("Only absolute http/https URLs are allowed")
    if parsed.username or parsed.password:
        raise ValueError("URLs with embedded credentials are not allowed")
    host = parsed.hostname
    try:
        addresses = {entry[4][0] for entry in socket.getaddrinfo(host, parsed.port or (443 if parsed.scheme == "https" else 80), type=socket.SOCK_STREAM)}
    except OSError as exc:
        raise ValueError(f"Could not resolve host {host}: {exc}") from exc
    if not addresses or any(not ipaddress.ip_address(address).is_global for address in addresses):
        raise ValueError("Private, loopback, link-local, and reserved network destinations are blocked")
    return urllib.parse.urlunsplit(parsed)


class _SafeRedirectHandler(urllib.request.HTTPRedirectHandler):
    def redirect_request(self, req, fp, code, msg, headers, newurl):
        return super().redirect_request(req, fp, code, msg, headers, _validate_public_url(newurl))


def _fetch_public_text(url: str) -> tuple[str, str, int, bool]:
    validated = _validate_public_url(url)
    request = urllib.request.Request(
        validated,
        headers={
            "User-Agent": "OpenCowork-CrewRuntime/1.0 (+local user initiated research)",
            "Accept": "text/html,application/json,text/plain,application/xml;q=0.9,*/*;q=0.2",
        },
    )
    opener = urllib.request.build_opener(_SafeRedirectHandler())
    with opener.open(request, timeout=20) as response:
        content_type = response.headers.get_content_type().lower()
        allowed = content_type.startswith("text/") or content_type in {"application/json", "application/xml", "application/xhtml+xml"}
        if not allowed:
            raise ValueError(f"Unsupported web content type: {content_type}")
        raw = response.read(MAX_WEB_BYTES + 1)
        truncated = len(raw) > MAX_WEB_BYTES
        if truncated:
            raw = raw[:MAX_WEB_BYTES]
        charset = response.headers.get_content_charset() or "utf-8"
        body = raw.decode(charset, errors="replace")
        final_url = _validate_public_url(response.geturl())
        return final_url, body, int(getattr(response, "status", 200)), truncated


class WebFetchInput(BaseModel):
    url: str = Field(description="Public http/https URL")


class WebFetchTool(BaseTool):
    name: str = "web_fetch"
    description: str = "Fetch readable text from a public web URL. Private and local network destinations are blocked."
    args_schema: type[BaseModel] = WebFetchInput

    def _run(self, url: str) -> str:
        def execute() -> str:
            final_url, body, status, truncated = _fetch_public_text(url)
            extractor = _TextExtractor()
            extractor.feed(body)
            text = extractor.text()
            limit_note = "\nDownload truncated safely after 1000000 bytes." if truncated else ""
            return f"URL: {final_url}\nHTTP: {status}{limit_note}\n\n{_truncate(text, 20_000)}"

        return _safe_result("web_fetch", execute)


class _DuckDuckGoParser(HTMLParser):
    def __init__(self) -> None:
        super().__init__(convert_charrefs=True)
        self.results: list[dict[str, str]] = []
        self.current: dict[str, str] | None = None
        self.capture_title = False
        self.capture_snippet = False

    @staticmethod
    def _classes(attrs: list[tuple[str, str | None]]) -> set[str]:
        value = next((value or "" for key, value in attrs if key == "class"), "")
        return set(value.split())

    def handle_starttag(self, tag: str, attrs: list[tuple[str, str | None]]) -> None:
        classes = self._classes(attrs)
        if tag == "a" and ("result__a" in classes or "result-link" in classes):
            href = next((value or "" for key, value in attrs if key == "href"), "")
            self.current = {"url": href, "title": "", "snippet": ""}
            self.capture_title = True
        elif self.current is not None and ("result__snippet" in classes or "result-snippet" in classes):
            self.capture_snippet = True

    def handle_endtag(self, tag: str) -> None:
        if tag == "a" and self.capture_title:
            self.capture_title = False
            if self.current is not None:
                self.results.append(self.current)
        if self.capture_snippet and tag in {"a", "div", "td"}:
            self.capture_snippet = False
            self.current = None

    def handle_data(self, data: str) -> None:
        if self.current is None:
            return
        if self.capture_title:
            self.current["title"] += data
        elif self.capture_snippet:
            self.current["snippet"] += data


def _unwrap_search_url(value: str) -> str:
    absolute = urllib.parse.urljoin("https://duckduckgo.com", html.unescape(value))
    parsed = urllib.parse.urlsplit(absolute)
    query = urllib.parse.parse_qs(parsed.query)
    redirected = query.get("uddg", [""])[0]
    return urllib.parse.unquote(redirected) if redirected else absolute


class _BingParser(HTMLParser):
    def __init__(self) -> None:
        super().__init__(convert_charrefs=True)
        self.results: list[dict[str, str]] = []
        self.current: dict[str, str] | None = None
        self.result_depth = 0
        self.heading_depth = 0
        self.capture_title = False
        self.capture_snippet = False

    @staticmethod
    def _classes(attrs: list[tuple[str, str | None]]) -> set[str]:
        value = next((value or "" for key, value in attrs if key == "class"), "")
        return set(value.split())

    def handle_starttag(self, tag: str, attrs: list[tuple[str, str | None]]) -> None:
        classes = self._classes(attrs)
        if self.current is None and tag == "li" and "b_algo" in classes:
            self.current = {"url": "", "title": "", "snippet": ""}
            self.result_depth = 1
            return
        if self.current is None:
            return
        self.result_depth += 1
        if tag == "h2":
            self.heading_depth = self.result_depth
        elif tag == "a" and self.heading_depth:
            self.current["url"] = next((value or "" for key, value in attrs if key == "href"), "")
            self.capture_title = True
        elif tag == "p" and self.current.get("url"):
            self.capture_snippet = True

    def handle_endtag(self, tag: str) -> None:
        if self.current is None:
            return
        if tag == "a" and self.capture_title:
            self.capture_title = False
        elif tag == "p" and self.capture_snippet:
            self.capture_snippet = False
        elif tag == "h2":
            self.heading_depth = 0
        self.result_depth -= 1
        if self.result_depth <= 0:
            if self.current.get("url") and self.current.get("title"):
                self.results.append(self.current)
            self.current = None
            self.result_depth = 0

    def handle_data(self, data: str) -> None:
        if self.current is None:
            return
        if self.capture_title:
            self.current["title"] += data
        elif self.capture_snippet:
            self.current["snippet"] += data


def _unwrap_bing_url(value: str) -> str:
    absolute = html.unescape(value)
    parsed = urllib.parse.urlsplit(absolute)
    if parsed.hostname and parsed.hostname.lower().endswith("bing.com"):
        encoded = urllib.parse.parse_qs(parsed.query).get("u", [""])[0]
        if encoded.startswith("a1"):
            try:
                payload = encoded[2:]
                payload += "=" * (-len(payload) % 4)
                decoded = base64.urlsafe_b64decode(payload.encode("ascii")).decode("utf-8")
                if decoded.startswith(("http://", "https://")):
                    return decoded
            except Exception:
                pass
    return absolute


class WebSearchInput(BaseModel):
    query: str = Field(description="Specific web search query")
    max_results: int = Field(default=6, ge=1, le=10)


class WebSearchTool(BaseTool):
    name: str = "web_search"
    description: str = "Search the public web and return titles, source URLs, and snippets for research."
    args_schema: type[BaseModel] = WebSearchInput

    def _run(self, query: str, max_results: int = 6) -> str:
        def execute() -> str:
            normalized = " ".join(str(query or "").split())
            if not normalized:
                raise ValueError("A non-empty search query is required")
            provider = "Bing"
            search_url = "https://www.bing.com/search?q=" + urllib.parse.quote_plus(normalized)
            _, body, _, _ = _fetch_public_text(search_url)
            parser: _BingParser | _DuckDuckGoParser = _BingParser()
            parser.feed(body)
            if not parser.results:
                provider = "DuckDuckGo"
                search_url = "https://html.duckduckgo.com/html/?q=" + urllib.parse.quote_plus(normalized)
                _, body, _, _ = _fetch_public_text(search_url)
                parser = _DuckDuckGoParser()
                parser.feed(body)
            rendered: list[str] = []
            seen: set[str] = set()
            for item in parser.results:
                raw_url = item.get("url", "")
                url = _unwrap_bing_url(raw_url) if provider == "Bing" else _unwrap_search_url(raw_url)
                if not url.startswith(("http://", "https://")) or url in seen:
                    continue
                seen.add(url)
                title = " ".join(item.get("title", "").split()) or url
                snippet = " ".join(item.get("snippet", "").split())
                rendered.append(f"{len(rendered) + 1}. {title}\nURL: {url}\nSnippet: {snippet or '(no snippet)'}")
                if len(rendered) >= max_results:
                    break
            if not rendered:
                raise RuntimeError("The search provider returned no parseable results; refine the query and retry")
            return f"Search query: {normalized}\nProvider: {provider}\nResults: {len(rendered)}\n\n" + "\n\n".join(rendered)

        return _safe_result("web_search", execute)


class BashInput(BaseModel):
    command: str = Field(description="Non-interactive command to run in the working directory")
    timeout_seconds: int = Field(default=60, ge=1, le=120)


def _subprocess_environment() -> dict[str, str]:
    blocked_python_variables = {
        "PYTHONHOME",
        "PYTHONPATH",
        "PYTHONEXECUTABLE",
        "__PYVENV_LAUNCHER__",
    }
    environment = {
        key: value
        for key, value in os.environ.items()
        if not re.search(r"(?i)(api[_-]?key|token|secret|password|credential)", key)
        and key.upper() not in blocked_python_variables
    }
    runtime_bin = str(Path(sys.executable).resolve().parent)
    existing_path = environment.get("PATH", "")
    environment["PATH"] = runtime_bin + (os.pathsep + existing_path if existing_path else "")
    return environment


class BashTool(BaseTool):
    name: str = "bash"
    description: str = "Run a bounded, non-interactive PowerShell command on Windows or POSIX shell command elsewhere, from the working directory."
    args_schema: type[BaseModel] = BashInput
    _root: Path = PrivateAttr()

    def __init__(self, root: Path) -> None:
        super().__init__()
        self._root = root

    def _run(self, command: str, timeout_seconds: int = 60) -> str:
        def execute() -> str:
            normalized = str(command or "").strip()
            if not normalized:
                raise ValueError("A non-empty command is required")
            destructive = re.compile(
                r"(?i)(git\s+(?:reset\s+--hard|clean\s+-[^\n]*f)|remove-item\b[^\n]*(?:-recurse|-force)|\brmdir\s+/s\b|\brd\s+/s\b|\bformat(?:\.com)?\b|\bshutdown\b|\bstop-computer\b)"
            )
            if destructive.search(normalized):
                raise PermissionError("Destructive shell commands are blocked by the Crew runtime")
            args = (
                ["powershell.exe", "-NoProfile", "-NonInteractive", "-Command", normalized]
                if os.name == "nt"
                else ["/bin/sh", "-lc", normalized]
            )
            completed = subprocess.run(
                args,
                cwd=self._root,
                env=_subprocess_environment(),
                capture_output=True,
                text=True,
                encoding="utf-8",
                errors="replace",
                timeout=timeout_seconds,
                check=False,
            )
            output = "\n".join(part for part in [completed.stdout.strip(), completed.stderr.strip()] if part)
            return f"Exit code: {completed.returncode}\n{_truncate(output or '(no output)', 20_000)}"

        return _safe_result("bash", execute)


class TodoInput(BaseModel):
    action: Literal["list", "add", "complete", "clear"]
    item: str = Field(default="", description="Todo text for add, or 1-based item number for complete")


class TodoTool(BaseTool):
    name: str = "todo"
    description: str = "Maintain a small in-memory checklist for this agent run."
    args_schema: type[BaseModel] = TodoInput
    _items: list[dict[str, Any]] = PrivateAttr(default_factory=list)

    def _run(self, action: str, item: str = "") -> str:
        def execute() -> str:
            if action == "add":
                text = " ".join(item.split())
                if not text:
                    raise ValueError("Todo text is required")
                self._items.append({"text": text, "done": False})
            elif action == "complete":
                index = int(item) - 1
                if index < 0 or index >= len(self._items):
                    raise ValueError("Todo number is out of range")
                self._items[index]["done"] = True
            elif action == "clear":
                self._items.clear()
            elif action != "list":
                raise ValueError(f"Unknown todo action: {action}")
            if not self._items:
                return "Todo list is empty."
            return "\n".join(f"{index}. [{'x' if value['done'] else ' '}] {value['text']}" for index, value in enumerate(self._items, 1))

        return _safe_result("todo", execute)


class OfficeWorkflowInput(BaseModel):
    output_path: str = Field(description="Output .pptx or .docx path inside the working directory")
    title: str = Field(description="Document or presentation title")
    sections_json: str = Field(
        description='JSON array of sections/slides. Each item may contain "title", "body", and "bullets" (string array).'
    )


def _parse_office_sections(value: str) -> list[dict[str, Any]]:
    parsed = json.loads(value)
    if isinstance(parsed, dict):
        parsed = parsed.get("sections") or parsed.get("slides") or []
    if not isinstance(parsed, list) or not parsed:
        raise ValueError("sections_json must contain a non-empty JSON array")
    sections: list[dict[str, Any]] = []
    for entry in parsed[:40]:
        if not isinstance(entry, dict):
            continue
        bullets = entry.get("bullets") or []
        if isinstance(bullets, str):
            bullets = [bullets]
        sections.append({
            "title": str(entry.get("title") or "Section").strip(),
            "body": str(entry.get("body") or "").strip(),
            "bullets": [str(item).strip() for item in bullets if str(item).strip()][:20],
        })
    if not sections:
        raise ValueError("sections_json did not contain any valid section objects")
    return sections


class OfficeWorkflowTool(BaseTool):
    name: str = "office_workflow"
    description: str = (
        "Create a real PowerPoint (.pptx) or Word (.docx) artifact. Call the tool directly with output_path, title, "
        "and sections_json. sections_json must be a JSON array such as "
        "[{\"title\":\"Evidence\",\"bullets\":[\"Verified fact\"]}]. Do not return a proposed tool call as text."
    )
    args_schema: type[BaseModel] = OfficeWorkflowInput
    _root: Path = PrivateAttr()

    def __init__(self, root: Path) -> None:
        super().__init__()
        self._root = root

    def _run(self, output_path: str, title: str, sections_json: str) -> str:
        def execute() -> str:
            target = _resolve_workspace_path(self._root, output_path, allow_root=False)
            suffix = target.suffix.lower()
            if suffix not in {".pptx", ".docx"}:
                raise ValueError("output_path must end in .pptx or .docx")
            sections = _parse_office_sections(sections_json)
            target.parent.mkdir(parents=True, exist_ok=True)
            if suffix == ".pptx":
                from pptx import Presentation  # type: ignore

                presentation = Presentation()
                title_slide = presentation.slides.add_slide(presentation.slide_layouts[0])
                title_slide.shapes.title.text = title.strip() or "Presentation"
                if len(title_slide.placeholders) > 1:
                    title_slide.placeholders[1].text = "Created by Open Cowork CrewAI"
                for section in sections:
                    slide = presentation.slides.add_slide(presentation.slide_layouts[1])
                    slide.shapes.title.text = section["title"]
                    frame = slide.placeholders[1].text_frame
                    frame.clear()
                    items = ([section["body"]] if section["body"] else []) + section["bullets"]
                    for index, item in enumerate(items or [""]):
                        paragraph = frame.paragraphs[0] if index == 0 else frame.add_paragraph()
                        paragraph.text = item
                        paragraph.level = 0
                presentation.save(target)
            else:
                from docx import Document  # type: ignore

                document = Document()
                document.add_heading(title.strip() or "Document", level=0)
                for section in sections:
                    document.add_heading(section["title"], level=1)
                    if section["body"]:
                        document.add_paragraph(section["body"])
                    for bullet in section["bullets"]:
                        document.add_paragraph(bullet, style="List Bullet")
                document.save(target)
            return f"Created {target} with {len(sections)} section(s) ({target.stat().st_size} bytes)."

        return _safe_result("office_workflow", execute)


TOOL_FACTORIES = {
    "read_file": lambda root: ReadFileTool(root),
    "edit_file": lambda root: EditFileTool(root),
    "create_directory": lambda root: CreateDirectoryTool(root),
    "move_path": lambda root: MovePathTool(root),
    "copy_path": lambda root: CopyPathTool(root),
    "glob": lambda root: GlobTool(root),
    "grep": lambda root: GrepTool(root),
    "web_fetch": lambda root: WebFetchTool(),
    "web_search": lambda root: WebSearchTool(),
    "bash": lambda root: BashTool(root),
    "todo": lambda root: TodoTool(),
    "office_workflow": lambda root: OfficeWorkflowTool(root),
}


def build_runtime_tools(request: dict, agent: dict) -> list[BaseTool]:
    agent_id = str(agent.get("id") or "").strip()
    requested = [_canonical_tool_id(value) for value in agent.get("tools") or [] if str(value).strip()]
    access = _agent_access(request, agent_id)
    allowed = {_canonical_tool_id(value) for value in access.get("allowedTools") or []}
    blocked = {_canonical_tool_id(value) for value in access.get("blockedTools") or []}
    root = _workspace_root(request)
    result: list[BaseTool] = []
    seen: set[str] = set()
    for tool_id in requested:
        if tool_id in seen or tool_id in blocked or tool_id not in allowed:
            continue
        factory = TOOL_FACTORIES.get(tool_id)
        if factory is None:
            continue
        result.append(factory(root))
        seen.add(tool_id)
    return result


def unavailable_runtime_tools(request: dict, agent: dict) -> list[str]:
    agent_id = str(agent.get("id") or "").strip()
    requested = {_canonical_tool_id(value) for value in agent.get("tools") or [] if str(value).strip()}
    access = _agent_access(request, agent_id)
    allowed = {_canonical_tool_id(value) for value in access.get("allowedTools") or []}
    return sorted(tool_id for tool_id in requested & allowed if tool_id not in TOOL_FACTORIES and tool_id not in {"delegate_task"})
